from pptx import Presentation

# Create presentation
prs = Presentation()

# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Improving Productivity Under Changing Client Requirements"
subtitle.text = "Guidance for Engineering Teams"

# Content sections
sections = {
    "Strengthen Requirement Management": [
        "Use a “Definition of Ready” for clear criteria",
        "Introduce change-control gates",
    ],
    "Improve Communication With Client": [
        "Hold weekly change review meetings",
        "Use impact matrices to show trade-offs",
    ],
    "Protect Engineer Focus Time": [
        "Use shielded sprints with locked requirements",
        "Provide deep-work blocks & reduce interruptions",
    ],
    "Adopt Agile Practices for Volatile Projects": [
        "Use Kanban for better flexibility",
        "Feature flags to reduce rework and risk",
    ],
    "Improve Internal Engineering Processes": [
        "Automate CI/CD, testing, and linting",
        "Use modular architecture to reduce impact of changes",
    ],
    "Manage Client Expectations": [
        "Document every change using CRs",
        "Communicate timeline impacts clearly",
    ],
    "Boost Team Morale": [
        "Give engineers decision-making input",
        "Recognize invisible work & manage burnout",
    ]
}

# Add slides
for title_text, bullets in sections.items():
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Body layout
    slide.shapes.title.text = title_text

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            body.text = bullet
        else:
            p = body.add_paragraph()
            p.text = bullet
            p.level = 0

# Save file
prs.save("Productivity_Changing_Requirements.pptx")
print("PPT generated: Productivity_Changing_Requirements.pptx")

